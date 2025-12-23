import cv2
import os
import sys
import numpy as np
import tempfile
import shutil
from pptx import Presentation
from pptx.util import Inches
from concurrent.futures import ThreadPoolExecutor, as_completed
import threading

# 创建线程锁用于保护打印输出
print_lock = threading.Lock()

def dhash(image, hash_size=8):
    """
    计算图像的差异哈希（dHash）
    用于快速比较两张图片的相似度，解决重复图片和渐进动画问题
    
    原理：将图像缩放到 (hash_size+1) x hash_size，然后比较相邻像素的亮度差异
    返回：64位整数哈希值
    """
    # 将图像缩放到 (hash_size+1) x hash_size
    resized = cv2.resize(image, (hash_size + 1, hash_size))
    # 计算水平方向的差异：如果右边像素比左边亮则为1，否则为0
    diff = resized[:, 1:] > resized[:, :-1]
    # 将布尔数组转换为整数哈希值
    return sum([2 ** i for i, v in enumerate(diff.flatten()) if v])

def hamming_distance(hash1, hash2):
    """
    计算两个哈希值的汉明距离（不同位的数量）
    用于判断两张图片的相似度
    """
    return bin(hash1 ^ hash2).count('1')

def calculate_robust_motion_score(frame1, frame2, threshold_pixel=25):
    """
    更鲁棒的运动检测：结合像素差异和结构相似度
    返回：运动分数（0-100），值越大表示变化越大
    """
    # 计算绝对差异
    diff = cv2.absdiff(frame1, frame2)
    
    # 使用自适应阈值，忽略微小变化
    _, thresh = cv2.threshold(diff, threshold_pixel, 255, cv2.THRESH_BINARY)
    
    # 计算变化区域占比
    motion_ratio = (np.count_nonzero(thresh) / thresh.size) * 100
    
    # 计算差异的均值（用于检测整体亮度变化）
    mean_diff = np.mean(diff)
    
    # 综合评分：结合变化区域占比和平均差异
    # 如果平均差异很大，即使变化区域小，也可能是整体亮度变化（翻页）
    if mean_diff > 10:  # 整体亮度变化超过阈值
        motion_ratio = max(motion_ratio, mean_diff / 2.55)  # 归一化到0-100范围
    
    return motion_ratio

def extract_slides_no_dedup(video_path, output_ppt_path, 
                            motion_threshold=1.5,     # 判定画面是否在动的阈值(%) - 调高以忽略噪点
                            min_duration=0.6,         # 画面需要静止多久才保存(秒) - 调低以提高灵敏度
                            hash_threshold=3,         # 哈希去重阈值：汉明距离小于此值认为是重复（更严格）
                            use_morphology=True,      # 是否使用形态学操作消除鼠标干扰
                            ignore_bottom_ratio=0.0,  # 忽略底部区域比例（0.0-1.0，如0.1表示忽略底部10%）
                            ignore_right_ratio=0.0,   # 忽略右侧区域比例（0.0-1.0，如0.1表示忽略右侧10%）
                            force_capture_interval=5.0  # 强制捕获间隔（秒）：长时间未捕获时强制检查
                            ):
    """
    优化版PPT提取函数 - "宁肯重复，绝不漏掉"策略
    
    核心优化：
    1. 状态机触发：检测到显著变化后，一旦进入相对稳定就立即捕获
    2. 鲁棒运动检测：改进的motion_score计算，更好地区分翻页和噪点
    3. 强制捕获机制：长时间未捕获时强制检查并保存
    4. 弱化去重：除非哈希值几乎完全一样，否则保留所有候选帧
    5. 形态学操作：通过开运算消除鼠标指针等小物体
    6. ROI机制：可选择性忽略底部（任务栏）和右侧（摄像头）区域
    """
    
    # 创建临时文件夹
    temp_dir = tempfile.mkdtemp()
    with print_lock:
        print(f"临时文件夹 created: {temp_dir}")
    
    try:
        cap = cv2.VideoCapture(video_path)
        if not cap.isOpened():
            with print_lock:
                print(f"错误：无法打开视频 {video_path}")
            return

        fps = cap.get(cv2.CAP_PROP_FPS)
        if fps <= 0: fps = 25 # 兜底防止读取不到fps
        
        # 采样处理间隔：每0.2秒检测一次画面状态（平衡性能与精度）
        process_interval = 0.2 
        frame_step = max(1, int(fps * process_interval))
        
        # 计算需要连续静止多少次循环才算"稳定"（用于快速捕获）
        frames_needed_for_stable = max(1, int(min_duration / process_interval))
        
        # 强制捕获：计算需要多少个周期后触发强制检查
        force_capture_frames = int(force_capture_interval / process_interval)

        slide_count = 0
        slide_images = []
        
        # 核心变量 - 新状态机
        last_processed_gray = None   # 上一检测帧（用来判断屏幕还在动吗？）
        stable_counter = 0           # 计数器：画面已经连续静止了多少个周期
        last_saved_hash = None       # 上一张已保存图片的哈希值（用于去重）
        pending_capture = False      # 待捕获标志：检测到显著变化后，等待稳定
        last_significant_change_frame = None  # 上一次显著变化时的帧（用于哈希对比）
        frames_since_last_capture = 0  # 距离上次捕获的帧数（用于强制捕获）
        last_capture_time = 0        # 上次捕获的时间（秒）
        
        with print_lock:
            print(f"开始处理视频: {os.path.basename(video_path)}")
            print(f"视频FPS: {fps:.2f}")
            print(f"策略: 检测到显著变化后，静止 {min_duration} 秒即立即捕获（最大化灵敏度）")
            print(f"motion_threshold: {motion_threshold}% (调高以忽略噪点)")
            if use_morphology:
                print(f"✓ 已启用形态学操作（消除鼠标干扰）")
            if ignore_bottom_ratio > 0 or ignore_right_ratio > 0:
                print(f"✓ ROI过滤: 底部{ignore_bottom_ratio*100:.0f}%, 右侧{ignore_right_ratio*100:.0f}%")

        while True:
            # 1. 快进读取（跳过中间帧，只读关键时间点）
            for _ in range(frame_step):
                cap.grab() 
            
            ret, frame = cap.read()
            if not ret:
                break

            # 2. 图像预处理
            gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
            
            # 【优化1：ROI机制】选择性忽略底部和右侧区域（任务栏、摄像头等）
            h, w = gray.shape
            roi_top = 0
            roi_bottom = int(h * (1 - ignore_bottom_ratio))
            roi_left = 0
            roi_right = int(w * (1 - ignore_right_ratio))
            gray_roi = gray[roi_top:roi_bottom, roi_left:roi_right]
            
            # 高斯模糊去噪，防止视频压缩噪点导致误判为"在动"
            gray_roi = cv2.GaussianBlur(gray_roi, (21, 21), 0)
            
            # 【优化2：形态学操作】消除鼠标指针、激光笔等小物体干扰
            # 开运算（先腐蚀后膨胀）可以去除小的噪点，只保留大的内容块
            if use_morphology:
                kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (5, 5))
                gray_roi = cv2.morphologyEx(gray_roi, cv2.MORPH_OPEN, kernel)
            
            # 将处理后的ROI区域放回原图（用于后续保存完整帧）
            gray_processed = gray.copy()
            gray_processed[roi_top:roi_bottom, roi_left:roi_right] = gray_roi

            # 初始化第一帧：保存第一帧作为基准，并标记为待捕获（如果视频开始就是静止的）
            if last_processed_gray is None:
                last_processed_gray = gray_processed
                pending_capture = True  # 标记第一帧为待捕获，如果接下来稳定就捕获
                continue

            # 3. 计算"动静"：使用更鲁棒的运动检测
            motion_score = calculate_robust_motion_score(last_processed_gray, gray_processed)
            
            # 获取当前时间
            current_time = cap.get(cv2.CAP_PROP_POS_MSEC) / 1000.0
            
            # 更新对比基准（使用处理后的图像）
            last_processed_gray = gray_processed
            frames_since_last_capture += 1

            # 4. 新状态机逻辑：检测显著变化 -> 待捕获 -> 稳定后立即捕获
            significant_change = motion_score >= motion_threshold * 2  # 显著变化的阈值（翻页动作）
            is_stable = motion_score < motion_threshold  # 相对稳定的阈值
            
            # 检测到显著变化（翻页动作）
            if significant_change:
                pending_capture = True  # 标记为待捕获状态
                last_significant_change_frame = gray_processed.copy()  # 保存变化时的帧
                stable_counter = 0  # 重置稳定计数器
                with print_lock:
                    print(f"【{os.path.basename(video_path)}】检测到显著变化 (motion={motion_score:.2f}%)，进入待捕获状态")
            elif is_stable:
                # 画面相对稳定
                if pending_capture:
                    # 在待捕获状态下，一旦稳定就立即捕获
                    stable_counter += 1
                    # 只要稳定了足够短的时间就捕获（不等待长时间）
                    if stable_counter >= frames_needed_for_stable:
                        # 立即捕获
                        curr_hash = dhash(gray_processed)
                        should_save = True
                        
                        # 弱化去重：除非哈希值几乎完全一样，否则保存
                        if last_saved_hash is not None:
                            hamming_dist = hamming_distance(curr_hash, last_saved_hash)
                            if hamming_dist < hash_threshold:  # 只有非常相似才跳过
                                should_save = False
                                with print_lock:
                                    print(f"【{os.path.basename(video_path)}】跳过几乎完全重复的帧（汉明距离={hamming_dist}）")
                        
                        if should_save:
                            img_path = os.path.join(temp_dir, f"slide_{slide_count:04d}.jpg")
                            cv2.imwrite(img_path, frame)
                            slide_images.append(img_path)
                            last_saved_hash = curr_hash
                            slide_count += 1
                            last_capture_time = current_time
                            frames_since_last_capture = 0
                            with print_lock:
                                print(f"【{os.path.basename(video_path)}】捕获第 {slide_count} 张 (于视频 {current_time:.1f}秒处, motion={motion_score:.2f}%)")
                        
                        # 重置状态
                        pending_capture = False
                        stable_counter = 0
                else:
                    # 不在待捕获状态，正常计数（用于常规稳定检测）
                    stable_counter += 1
            else:
                # 画面在轻微变动（可能是噪点或微弱动画）
                # 如果已经稳定了一段时间，轻微变动不影响
                if stable_counter > 0:
                    stable_counter = max(0, stable_counter - 1)  # 轻微减少，不完全清零
                else:
                    stable_counter = 0

            # 5. 强制捕获机制：如果长时间未捕获，强制检查并保存
            # 包括两种情况：1) 长时间未捕获；2) 待捕获状态但一直未稳定
            if frames_since_last_capture >= force_capture_frames:
                was_pending = pending_capture  # 记录是否在待捕获状态
                curr_hash = dhash(gray_processed)
                should_save = True
                
                # 与上一张对比，只有几乎完全一样才跳过
                if last_saved_hash is not None:
                    hamming_dist = hamming_distance(curr_hash, last_saved_hash)
                    if hamming_dist < hash_threshold:
                        should_save = False
                        with print_lock:
                            print(f"【{os.path.basename(video_path)}】强制检查：跳过重复帧（汉明距离={hamming_dist}）")
                
                if should_save:
                    img_path = os.path.join(temp_dir, f"slide_{slide_count:04d}.jpg")
                    cv2.imwrite(img_path, frame)
                    slide_images.append(img_path)
                    last_saved_hash = curr_hash
                    slide_count += 1
                    last_capture_time = current_time
                    frames_since_last_capture = 0
                    pending_capture = False  # 重置待捕获状态
                    with print_lock:
                        reason = "待捕获但一直未稳定" if was_pending else f"距离上次捕获 {force_capture_interval}秒"
                        print(f"【{os.path.basename(video_path)}】强制捕获第 {slide_count} 张 (于视频 {current_time:.1f}秒处, {reason})")

        cap.release()
        
        # 6. 【弱化后处理去重】只删除几乎完全相同的图片，保留所有候选帧
        # 策略：除非哈希值几乎完全一样（汉明距离 < hash_threshold），否则全部保留
        original_count = len(slide_images)
        if len(slide_images) > 1:
            filtered_images = [slide_images[0]]  # 保留第一张
            for i in range(1, len(slide_images)):
                prev_img = cv2.imread(filtered_images[-1], cv2.IMREAD_GRAYSCALE)
                curr_img = cv2.imread(slide_images[i], cv2.IMREAD_GRAYSCALE)
                
                if prev_img is None or curr_img is None:
                    filtered_images.append(slide_images[i])
                    continue
                
                # 计算两张图片的哈希值
                prev_hash = dhash(prev_img)
                curr_hash = dhash(curr_img)
                hamming_dist = hamming_distance(prev_hash, curr_hash)
                
                # 弱化去重：只有哈希值几乎完全一样才删除（非常严格的阈值）
                if hamming_dist < hash_threshold:
                    # 几乎完全重复，跳过
                    with print_lock:
                        print(f"【{os.path.basename(video_path)}】后处理：删除几乎完全重复的帧（汉明距离={hamming_dist}）")
                else:
                    # 有任何差异都保留（宁肯重复，绝不漏掉）
                    filtered_images.append(slide_images[i])
            
            slide_images = filtered_images
            removed_count = original_count - len(slide_images)
            if removed_count > 0:
                with print_lock:
                    print(f"【{os.path.basename(video_path)}】后处理去重完成：保留 {len(slide_images)} 张，删除 {removed_count} 张几乎完全重复的帧")
            else:
                with print_lock:
                    print(f"【{os.path.basename(video_path)}】后处理：所有 {len(slide_images)} 张图片均保留（无完全重复）")
        
        # 7. 生成PPT
        if not slide_images:
            with print_lock:
                print(f"【{os.path.basename(video_path)}】未提取到任何图片，请检查阈值设置。")
            return

        with print_lock:
            print(f"【{os.path.basename(video_path)}】共提取 {len(slide_images)} 张图片，开始生成PPT...")
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        for i, img_path in enumerate(slide_images):
            blank_slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 读取图片并计算自适应尺寸
            img = cv2.imread(img_path)
            if img is None: continue
            
            img_height, img_width = img.shape[:2]
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            img_ratio = img_width / img_height
            slide_ratio = slide_width / slide_height
            
            if img_ratio > slide_ratio:
                width = slide_width
                height = int(slide_width / img_ratio)
                left = 0
                top = (slide_height - height) // 2
            else:
                height = slide_height
                width = int(slide_height * img_ratio)
                top = 0
                left = (slide_width - width) // 2
            
            slide.shapes.add_picture(img_path, left, top, width, height)
            if (i+1) % 10 == 0:
                with print_lock:
                    print(f"【{os.path.basename(video_path)}】已处理 {i+1}/{len(slide_images)} 张...")
        
        prs.save(output_ppt_path)
        with print_lock:
            print(f"【{os.path.basename(video_path)}】成功！PPT已保存至: {output_ppt_path}")
        
    finally:
        # 清理临时文件
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir, ignore_errors=True)
            with print_lock:
                print(f"【{os.path.basename(video_path)}】临时文件清理完成")

def extract_slides_dual_stream(video_path, output_base_dir,
                               # Group A 参数（高精度全量组）
                               group_a_motion_threshold=1.5,  # 运动阈值：检测显著变化（翻页动作）
                               group_a_min_duration=0.6,      # 静止时长（秒）：画面稳定后保存
                               group_a_quality=100,          # 原始画质保存（JPEG质量 1-100）
                               # Group B 参数（低精度去重组）
                               group_b_motion_threshold=1.5,  # 运动阈值：低于此值认为内容未变化
                               group_b_hash_threshold=5,      # 深度去重阈值：汉明距离小于此值认为重复
                               group_b_scale_factor=0.7,      # 分辨率缩放因子（0.7表示缩放到70%）
                               group_b_quality=75,            # JPEG压缩质量（降低文件体积）
                               # 通用参数
                               use_morphology=True,
                               ignore_bottom_ratio=0.0,
                               ignore_right_ratio=0.0):
    """
    双流（Dual-Stream）采集模块
    
    Group A: 高精度全量组 (Full-Fidelity Archive)
    - 智能触发：检测到显著变化后，一旦稳定就保存
    - 不去重，原始画质保存（保留所有捕获的帧）
    - 直接生成PPT文件
    
    Group B: 低精度去重组 (Optimized Summary)
    - 深度去重：即使时间戳在变，如果核心内容没变化就跳过
    - 降低分辨率或增加压缩比，减小文件体积
    - 直接生成PPT文件
    """
    # 创建临时文件夹用于存储图片
    temp_dir = tempfile.mkdtemp()
    group_a_temp_dir = os.path.join(temp_dir, "group_a")
    group_b_temp_dir = os.path.join(temp_dir, "group_b")
    os.makedirs(group_a_temp_dir, exist_ok=True)
    os.makedirs(group_b_temp_dir, exist_ok=True)
    
    # 输出PPT文件路径
    video_name = os.path.splitext(os.path.basename(video_path))[0]
    group_a_ppt = os.path.join(output_base_dir, f"{video_name}_高精度.pptx")
    group_b_ppt = os.path.join(output_base_dir, f"{video_name}_低精度.pptx")
    
    with print_lock:
        print(f"\n{'='*60}")
        print(f"开始双流采集: {os.path.basename(video_path)}")
        print(f"Group A (高精度全量): 将生成 -> {group_a_ppt}")
        print(f"Group B (低精度去重): 将生成 -> {group_b_ppt}")
        print(f"{'='*60}")
    
    try:
        cap = cv2.VideoCapture(video_path)
        if not cap.isOpened():
            with print_lock:
                print(f"错误：无法打开视频 {video_path}")
            return
        
        fps = cap.get(cv2.CAP_PROP_FPS)
        if fps <= 0: fps = 25
        
        # 采样处理间隔：每0.2秒检测一次画面状态（平衡性能与精度）
        process_interval = 0.2
        frame_step = max(1, int(fps * process_interval))
        
        # Group A: 计算需要连续静止多少次循环才算"稳定"
        frames_needed_for_stable = max(1, int(group_a_min_duration / process_interval))
        
        # Group B: 使用相同的间隔进行检测
        group_b_interval = 0.2
        group_b_frame_step = max(1, int(fps * group_b_interval))
        
        # Group A 变量（智能触发状态机）
        group_a_images = []  # 存储图片路径列表
        group_a_count = 0
        group_a_last_gray = None
        group_a_stable_counter = 0
        group_a_pending_capture = False
        
        # Group B 变量
        group_b_images = []  # 存储图片路径列表
        group_b_count = 0
        last_group_b_gray = None
        last_group_b_hash = None
        
        frame_idx = 0
        
        while True:
            # 快进读取（跳过中间帧，只读关键时间点）
            for _ in range(frame_step):
                cap.grab()
            
            ret, frame = cap.read()
            if not ret:
                break
            
            current_time_ms = cap.get(cv2.CAP_PROP_POS_MSEC)
            current_time_sec = current_time_ms / 1000.0
            
            # 图像预处理（ROI + 高斯模糊 + 形态学）
            gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
            h, w = gray.shape
            roi_top = 0
            roi_bottom = int(h * (1 - ignore_bottom_ratio))
            roi_left = 0
            roi_right = int(w * (1 - ignore_right_ratio))
            gray_roi = gray[roi_top:roi_bottom, roi_left:roi_right]
            gray_roi = cv2.GaussianBlur(gray_roi, (21, 21), 0)
            
            if use_morphology:
                kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (5, 5))
                gray_roi = cv2.morphologyEx(gray_roi, cv2.MORPH_OPEN, kernel)
            
            gray_processed = gray.copy()
            gray_processed[roi_top:roi_bottom, roi_left:roi_right] = gray_roi
            
            # ========== Group A: 高精度全量组 ==========
            # 智能触发：检测到显著变化后，一旦稳定就保存（不去重）
            if group_a_last_gray is None:
                # 第一帧：直接保存并初始化
                img_path = os.path.join(group_a_temp_dir, f"frame_{group_a_count:06d}.jpg")
                cv2.imwrite(img_path, frame, [cv2.IMWRITE_JPEG_QUALITY, group_a_quality])
                group_a_images.append(img_path)
                group_a_count += 1
                
                group_a_last_gray = gray_processed
                group_a_pending_capture = True
                group_a_stable_counter = 0
            else:
                # 计算运动分数
                motion_score = calculate_robust_motion_score(group_a_last_gray, gray_processed)
                
                # 检测显著变化（翻页动作）
                significant_change = motion_score >= group_a_motion_threshold * 2
                is_stable = motion_score < group_a_motion_threshold
                
                if significant_change:
                    # 检测到显著变化，进入待捕获状态
                    group_a_pending_capture = True
                    group_a_stable_counter = 0
                elif is_stable:
                    # 画面相对稳定
                    if group_a_pending_capture:
                        group_a_stable_counter += 1
                        # 一旦稳定了足够时间就捕获
                        if group_a_stable_counter >= frames_needed_for_stable:
                            # 保存（不去重，保留所有捕获的帧）
                            img_path = os.path.join(group_a_temp_dir, f"frame_{group_a_count:06d}.jpg")
                            cv2.imwrite(img_path, frame, [cv2.IMWRITE_JPEG_QUALITY, group_a_quality])
                            group_a_images.append(img_path)
                            group_a_count += 1
                            
                            # 更新基准并重置状态
                            group_a_last_gray = gray_processed
                            group_a_pending_capture = False
                            group_a_stable_counter = 0
                            
                            if group_a_count % 10 == 0:
                                with print_lock:
                                    print(f"【Group A】已捕获 {group_a_count} 张 (时间: {current_time_sec:.2f}秒)")
                else:
                    # 画面在轻微变动，重置稳定计数器
                    if group_a_stable_counter > 0:
                        group_a_stable_counter = max(0, group_a_stable_counter - 1)
                    else:
                        group_a_stable_counter = 0
            
            # ========== Group B: 低精度去重组 ==========
            # 只在特定间隔检测（节省计算）
            if frame_idx % group_b_frame_step == 0:
                if last_group_b_gray is None:
                    # 第一帧：保存并初始化
                    last_group_b_gray = gray_processed
                    last_group_b_hash = dhash(gray_processed)
                    
                    # 保存第一帧（低精度）
                    img_path = os.path.join(group_b_temp_dir, f"frame_{group_b_count:06d}.jpg")
                    frame_resized = cv2.resize(frame, 
                                             (int(w * group_b_scale_factor), 
                                              int(h * group_b_scale_factor)))
                    cv2.imwrite(img_path, frame_resized, 
                              [cv2.IMWRITE_JPEG_QUALITY, group_b_quality])
                    group_b_images.append(img_path)
                    group_b_count += 1
                else:
                    # 计算运动分数和哈希值
                    motion_score = calculate_robust_motion_score(last_group_b_gray, gray_processed)
                    curr_hash = dhash(gray_processed)
                    
                    # 深度去重判断
                    hamming_dist = hamming_distance(curr_hash, last_group_b_hash) if last_group_b_hash else 999
                    
                    # 判断是否保存：内容有实质变化（运动分数高 或 哈希差异大）
                    content_changed = (motion_score >= group_b_motion_threshold or 
                                     hamming_dist >= group_b_hash_threshold)
                    
                    if content_changed:
                        # 内容有实质变化，保存
                        img_path = os.path.join(group_b_temp_dir, f"frame_{group_b_count:06d}.jpg")
                        
                        # 降低分辨率并压缩
                        frame_resized = cv2.resize(frame, 
                                                 (int(w * group_b_scale_factor), 
                                                  int(h * group_b_scale_factor)))
                        cv2.imwrite(img_path, frame_resized, 
                                  [cv2.IMWRITE_JPEG_QUALITY, group_b_quality])
                        group_b_images.append(img_path)
                        
                        # 更新基准
                        last_group_b_gray = gray_processed
                        last_group_b_hash = curr_hash
                        group_b_count += 1
                        
                        with print_lock:
                            print(f"【Group B】捕获第 {group_b_count} 张 (时间: {current_time_sec:.1f}秒, "
                                f"motion={motion_score:.2f}%, hamming={hamming_dist})")
            
            frame_idx += 1
        
        cap.release()
        
        # 生成PPT文件
        def create_ppt_from_images(images, output_ppt_path, group_name):
            """从图片列表生成PPT"""
            if not images:
                with print_lock:
                    print(f"【{group_name}】未提取到任何图片，跳过PPT生成")
                return
            
            with print_lock:
                print(f"【{group_name}】共提取 {len(images)} 张图片，开始生成PPT...")
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            for i, img_path in enumerate(images):
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # 读取图片并计算自适应尺寸
                img = cv2.imread(img_path)
                if img is None:
                    continue
                
                img_height, img_width = img.shape[:2]
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                img_ratio = img_width / img_height
                slide_ratio = slide_width / slide_height
                
                if img_ratio > slide_ratio:
                    width = slide_width
                    height = int(slide_width / img_ratio)
                    left = 0
                    top = (slide_height - height) // 2
                else:
                    height = slide_height
                    width = int(slide_height * img_ratio)
                    top = 0
                    left = (slide_width - width) // 2
                
                slide.shapes.add_picture(img_path, left, top, width, height)
                if (i+1) % 10 == 0:
                    with print_lock:
                        print(f"【{group_name}】已处理 {i+1}/{len(images)} 张...")
            
            prs.save(output_ppt_path)
            with print_lock:
                print(f"【{group_name}】PPT已保存: {output_ppt_path}")
        
        # 生成Group A PPT
        create_ppt_from_images(group_a_images, group_a_ppt, "Group A (高精度全量)")
        
        # 生成Group B PPT
        create_ppt_from_images(group_b_images, group_b_ppt, "Group B (低精度去重)")
        
        with print_lock:
            print(f"\n{'='*60}")
            print(f"双流采集完成: {os.path.basename(video_path)}")
            print(f"Group A (高精度全量): {group_a_count} 张 -> {group_a_ppt}")
            print(f"Group B (低精度去重): {group_b_count} 张 -> {group_b_ppt}")
            print(f"{'='*60}\n")
    
    except Exception as e:
        with print_lock:
            print(f"错误：处理视频 {video_path} 时发生异常: {str(e)}")
            import traceback
            traceback.print_exc()
    finally:
        # 清理临时文件
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir, ignore_errors=True)
            with print_lock:
                print(f"临时文件清理完成")

# ================= 配置区域 =================
# 支持的视频格式
VIDEO_EXTENSIONS = ['.mp4', '.avi', '.mov', '.mkv', '.flv', '.wmv', '.m4v']

# 输出文件夹名称
OUTPUT_FOLDER = "output_ppts"

# 获取当前脚本所在目录
current_dir = os.path.dirname(os.path.abspath(__file__)) if __file__ else os.getcwd()

# 视频文件夹名称
VIDEO_FOLDER = "视频"

# 视频文件夹路径
video_dir = os.path.join(current_dir, VIDEO_FOLDER)

# 检查视频文件夹是否存在
if not os.path.exists(video_dir):
    print(f"错误：未找到视频文件夹 '{VIDEO_FOLDER}'")
    print(f"请确保在 {current_dir} 目录下存在 '{VIDEO_FOLDER}' 文件夹")
    sys.exit(1)

# 创建输出文件夹
output_dir = os.path.join(current_dir, OUTPUT_FOLDER)
if not os.path.exists(output_dir):
    os.makedirs(output_dir)
    print(f"已创建输出文件夹: {output_dir}")

# 查找所有视频文件（在视频文件夹中）
video_files = []
for file in os.listdir(video_dir):
    file_lower = file.lower()
    if any(file_lower.endswith(ext) for ext in VIDEO_EXTENSIONS):
        video_path = os.path.join(video_dir, file)
        video_files.append(video_path)

if not video_files:
    print(f"错误：在目录 {video_dir} 中未找到任何视频文件")
    print(f"支持的格式: {', '.join(VIDEO_EXTENSIONS)}")
else:
    print(f"找到 {len(video_files)} 个视频文件，开始双流采集处理...\n")
    
    # 处理单个视频的包装函数
    def process_video(video_file, idx, total):
        video_name = os.path.basename(video_file)
        video_name_no_ext = os.path.splitext(video_name)[0]
        
        with print_lock:
            print(f"\n{'='*60}")
            print(f"线程开始处理 [{idx}/{total}]: {video_name}")
            print(f"{'='*60}")
        
        try:
            # 双流采集模式
            extract_slides_dual_stream(
                video_file,
                output_dir,
                # Group A 参数（高精度全量组）
                group_a_motion_threshold=1.5,  # 运动阈值：检测显著变化
                group_a_min_duration=0.6,      # 静止时长：画面稳定后保存（0.5-0.8推荐）
                group_a_quality=100,            # 原始画质保存
                # Group B 参数（低精度去重组）
                group_b_motion_threshold=1.5,  # 运动阈值
                group_b_hash_threshold=5,       # 深度去重阈值
                group_b_scale_factor=0.7,       # 分辨率缩放（70%）
                group_b_quality=75,              # JPEG压缩质量
                # 通用参数
                use_morphology=True,
                ignore_bottom_ratio=0.0,
                ignore_right_ratio=0.0
            )
            
            with print_lock:
                print(f"✓ 完成 [{idx}/{total}]: {video_name}")
            return True, video_name
        except Exception as e:
            with print_lock:
                print(f"✗ 失败 [{idx}/{total}]: {video_name} - 错误: {str(e)}")
            return False, video_name
    
    # 使用多线程处理（线程数设置为CPU核心数，但不超过视频文件数量）
    import multiprocessing
    max_workers = min(multiprocessing.cpu_count(), len(video_files))
    print(f"使用 {max_workers} 个线程并行处理...\n")
    
    # 创建任务列表
    tasks = []
    for idx, video_file in enumerate(video_files, 1):
        tasks.append((video_file, idx, len(video_files)))
    
    # 使用线程池执行
    completed = 0
    failed = 0
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        # 提交所有任务
        future_to_video = {
            executor.submit(process_video, video_file, idx, len(video_files)): (idx, video_file)
            for idx, video_file in enumerate(video_files, 1)
        }
        
        # 等待所有任务完成
        for future in as_completed(future_to_video):
            idx, video_file = future_to_video[future]
            try:
                success, video_name = future.result()
                if success:
                    completed += 1
                else:
                    failed += 1
            except Exception as e:
                with print_lock:
                    print(f"✗ 异常 [{idx}/{len(video_files)}]: {os.path.basename(video_file)} - {str(e)}")
                failed += 1
    
    print(f"\n{'='*60}")
    print(f"批量处理完成！")
    print(f"成功: {completed} 个，失败: {failed} 个")
    print(f"所有PPT已保存至: {output_dir}")
    print(f"{'='*60}")